#!/usr/bin/env python3
"""
Generate a modern, professional PowerPoint presentation for Company Directory
SharePoint Hackathon 2026 - Intro Segment (Segment 1)

Design: Fluent 2-inspired, dark theme with Microsoft blue accents
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette (Fluent 2 / Microsoft brand) ─────────────────────────────
DARK_BG = RGBColor(0x1B, 0x1B, 0x1F)           # Dark surface
DARK_BG_ALT = RGBColor(0x24, 0x24, 0x2A)       # Slightly lighter dark
BRAND_BLUE = RGBColor(0x00, 0x78, 0xD4)        # Microsoft blue
BRAND_BLUE_LIGHT = RGBColor(0x47, 0xA1, 0xF5)  # Light accent blue
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xB8)
MID_GRAY = RGBColor(0x6E, 0x6E, 0x78)
SUCCESS_GREEN = RGBColor(0x0F, 0x7B, 0x0F)     # Fluent green
ACCENT_PURPLE = RGBColor(0x88, 0x64, 0xD8)     # Accent purple

# ─── Paths ───────────────────────────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
IMAGES_DIR = os.path.join(SCRIPT_DIR, "images")
OVERVIEW_IMG = os.path.join(IMAGES_DIR, "company-directory-overview.png")
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "Company-Directory-Hackathon-2026.pptx")

# ─── Helpers ─────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, alpha=None):
    """Add a rectangle shape with optional transparency."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no border
    if alpha is not None:
        # Set transparency via XML manipulation
        from lxml import etree
        solidFill = shape.fill._fill
        srgb = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is None:
            srgb = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr')
        if srgb is not None:
            nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            alpha_elem = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))  # val is in 1000ths of a percent
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI", spacing_after=Pt(0)):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = spacing_after
    return txBox


def add_multi_text(slide, left, top, width, height, lines, default_size=16,
                   default_color=LIGHT_GRAY, font_name="Segoe UI"):
    """Add a text box with multiple styled paragraphs.
    lines: list of dicts with keys: text, size, color, bold, alignment, spacing_after
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line.get("text", "")
        p.font.size = Pt(line.get("size", default_size))
        p.font.color.rgb = line.get("color", default_color)
        p.font.bold = line.get("bold", False)
        p.font.name = font_name
        p.alignment = line.get("alignment", PP_ALIGN.LEFT)
        p.space_after = line.get("spacing_after", Pt(4))

    return txBox


def add_accent_line(slide, left, top, width, color=BRAND_BLUE, height=Pt(3)):
    """Add a thin horizontal accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_pill_badge(slide, left, top, text, fill_color=BRAND_BLUE, text_color=WHITE, 
                   font_size=11, width=None):
    """Add a rounded-rectangle pill badge with text."""
    if width is None:
        width = Inches(len(text) * 0.09 + 0.4)
    height = Inches(0.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.35  # more rounded

    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "Segoe UI Semibold"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    # Vertical center
    from pptx.oxml.ns import qn
    txBody = tf._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    bodyPr.set('lIns', '72000')
    bodyPr.set('rIns', '72000')
    bodyPr.set('tIns', '0')
    bodyPr.set('bIns', '0')
    return shape


# ─── Slide Builders ──────────────────────────────────────────────────────────

def build_slide_1_title(prs):
    """Slide 1: Title Card - Company Directory
    Dark background, large title, subtitle, event badge
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide, DARK_BG)

    # Decorative accent bar at top
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), BRAND_BLUE)

    # Left decorative vertical line
    add_shape_rect(slide, Inches(0.8), Inches(1.5), Pt(3), Inches(4.0), BRAND_BLUE)

    # Main title
    add_text_box(slide, Inches(1.2), Inches(1.8), Inches(8), Inches(1.2),
                 "Company Directory", font_size=54, color=WHITE, bold=True,
                 font_name="Segoe UI Light")

    # Subtitle
    add_multi_text(slide, Inches(1.2), Inches(3.1), Inches(8), Inches(1.5), [
        {"text": "AI-Powered People Directory for Microsoft 365",
         "size": 24, "color": BRAND_BLUE_LIGHT, "bold": False,
         "spacing_after": Pt(16)},
        {"text": "Natural language search  •  Interactive org chart  •  Fluent UI v9",
         "size": 16, "color": LIGHT_GRAY, "bold": False,
         "spacing_after": Pt(24)},
    ])

    # Author info
    add_text_box(slide, Inches(1.2), Inches(4.4), Inches(6), Inches(0.5),
                 "João Mendes", font_size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.2), Inches(4.85), Inches(6), Inches(0.5),
                 "SharePoint & Microsoft 365 Developer", font_size=14, color=MID_GRAY)

    # Event badge
    add_pill_badge(slide, Inches(1.2), Inches(5.6), "SHAREPOINT HACKATHON 2026",
                   fill_color=BRAND_BLUE, width=Inches(3.2), font_size=12)

    # Right side: subtle background decoration - large faded circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(0.5),
                                     Inches(5.5), Inches(5.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x28)
    circle.line.fill.background()
    # Send to back
    from pptx.oxml.ns import qn
    sp = circle._element
    spTree = sp.getparent()
    spTree.remove(sp)
    spTree.insert(2, sp)  # insert near beginning to push behind text

    # Inner decorative circle
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.3), Inches(1.3),
                                      Inches(3.9), Inches(3.9))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(0x28, 0x28, 0x30)
    circle2.line.fill.background()
    sp2 = circle2._element
    spTree2 = sp2.getparent()
    spTree2.remove(sp2)
    spTree2.insert(3, sp2)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        "SLIDE 1 — TITLE (show for ~10 seconds)\n"
        "\n"
        "SAY:\n"
        "\"Hi everyone! I'm João Mendes, and this is Company Directory — "
        "an AI-powered people directory built entirely on the Microsoft 365 platform.\"\n"
        "\n"
        "TIPS:\n"
        "• Speak slowly, smile — this is the first impression.\n"
        "• Let the slide sit for 2-3 seconds before you start talking.\n"
        "• Keep your energy up — the viewer decides in the first 5 seconds whether to keep watching.\n"
        "• Don't read the slide — the audience can see it. Add your personality."
    )

    return slide


def build_slide_2_overview(prs):
    """Slide 2: Visual Overview
    Screenshot in center/right with dark overlay, key highlights on left
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Top accent bar
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.04), BRAND_BLUE)

    # Section label
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "OVERVIEW", font_size=12, color=BRAND_BLUE, bold=True,
                 font_name="Segoe UI Semibold")

    # Left column: key features
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(4.5), Inches(0.7),
                 "Your People,\nOne Search Away", font_size=32, color=WHITE, bold=True)

    features = [
        {"text": "", "size": 10, "color": DARK_BG, "spacing_after": Pt(12)},
        {"text": "▸  Browse employees in Grid or List view",
         "size": 15, "color": LIGHT_GRAY, "spacing_after": Pt(10)},
        {"text": "▸  AI search — ask in plain English",
         "size": 15, "color": LIGHT_GRAY, "spacing_after": Pt(10)},
        {"text": "▸  Interactive organization chart",
         "size": 15, "color": LIGHT_GRAY, "spacing_after": Pt(10)},
        {"text": "▸  Quick actions: chat, email, call, video",
         "size": 15, "color": LIGHT_GRAY, "spacing_after": Pt(10)},
        {"text": "▸  Works in SharePoint, Teams, Outlook & Office",
         "size": 15, "color": LIGHT_GRAY, "spacing_after": Pt(10)},
        {"text": "▸  Infinite scroll with IndexedDB caching",
         "size": 15, "color": LIGHT_GRAY, "spacing_after": Pt(10)},
    ]
    add_multi_text(slide, Inches(0.8), Inches(2.2), Inches(4.8), Inches(3.5), features)

    # Right column: screenshot or placeholder
    # Add a rounded container for the image
    img_left = Inches(5.8)
    img_top = Inches(0.8)
    img_width = Inches(7.0)
    img_height = Inches(5.8)

    # Dark frame behind image
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    img_left - Pt(4), img_top - Pt(4),
                                    img_width + Pt(8), img_height + Pt(8))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(0x30, 0x30, 0x38)
    frame.line.fill.background()
    frame.adjustments[0] = 0.02

    if os.path.exists(OVERVIEW_IMG):
        pic = slide.shapes.add_picture(OVERVIEW_IMG, img_left, img_top,
                                        img_width, img_height)
    else:
        # Placeholder
        placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              img_left, img_top, img_width, img_height)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = DARK_BG_ALT
        placeholder.line.color.rgb = MID_GRAY
        placeholder.line.width = Pt(1)
        add_text_box(slide, img_left + Inches(1.5), img_top + Inches(2.2),
                     Inches(4), Inches(1),
                     "[ Insert Grid View Screenshot ]",
                     font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        "SLIDE 2 — OVERVIEW (show for ~15 seconds)\n"
        "\n"
        "SAY:\n"
        "\"Company Directory solves a common problem — finding the right person in your organization, fast.\n"
        "\n"
        "You can browse employees in a visual Grid View or a compact List View, "
        "search using plain English powered by Azure OpenAI, explore the organization chart interactively, "
        "and take quick actions like starting a Teams chat or email — all from one web part.\n"
        "\n"
        "It works across SharePoint, Teams, Outlook, and Microsoft 365 — adapting to each host's theme.\"\n"
        "\n"
        "TIPS:\n"
        "• Point to the screenshot on the right if visible — guide the viewer's eye.\n"
        "• Read the bullet points only as prompts, don't read them word-for-word.\n"
        "• This is the \"why\" slide — make the audience care about the problem you're solving."
    )

    return slide


def build_slide_3_grid_view(prs):
    """Slide 3: Grid View Focus
    Full-width screenshot area with caption overlay
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Top accent bar
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.04), BRAND_BLUE)

    # Section label top-left
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.4),
                 "GRID VIEW", font_size=12, color=BRAND_BLUE, bold=True,
                 font_name="Segoe UI Semibold")

    # Left column: descriptive text
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(4.5), Inches(0.7),
                 "Visual Employee\nDirectory", font_size=30, color=WHITE, bold=True)

    grid_features = [
        {"text": "", "size": 8, "color": DARK_BG, "spacing_after": Pt(10)},
        {"text": "Each card displays the employee's photo, name, job title, department, and office — pulled live from Microsoft Graph.",
         "size": 15, "color": LIGHT_GRAY, "spacing_after": Pt(14)},
        {"text": "▸  Infinite scroll — loads 100 users at a time",
         "size": 14, "color": LIGHT_GRAY, "spacing_after": Pt(8)},
        {"text": "▸  No pagination — just scroll for more",
         "size": 14, "color": LIGHT_GRAY, "spacing_after": Pt(8)},
        {"text": "▸  IndexedDB caching for instant revisits",
         "size": 14, "color": LIGHT_GRAY, "spacing_after": Pt(8)},
        {"text": "▸  Responsive layout adapts to any screen",
         "size": 14, "color": LIGHT_GRAY, "spacing_after": Pt(8)},
        {"text": "▸  Click a card to view full profile & org chart",
         "size": 14, "color": LIGHT_GRAY, "spacing_after": Pt(8)},
    ]
    add_multi_text(slide, Inches(0.8), Inches(2.1), Inches(4.5), Inches(3.8), grid_features)

    # Right column: Grid View screenshot placeholder
    img_left = Inches(5.8)
    img_top = Inches(0.9)
    img_width = Inches(7.0)
    img_height = Inches(5.5)

    # Subtle frame
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    img_left, img_top, img_width, img_height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = DARK_BG_ALT
    frame.line.color.rgb = RGBColor(0x3A, 0x3A, 0x42)
    frame.line.width = Pt(1)
    frame.adjustments[0] = 0.015

    # Placeholder text
    add_text_box(slide, img_left + Inches(1.5), Inches(3.2), Inches(4), Inches(1),
                 "[ Insert Grid View Screenshot Here ]",
                 font_size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom caption bar
    add_shape_rect(slide, Inches(0), Inches(6.7), Inches(13.333), Inches(0.8), DARK_BG_ALT)

    add_multi_text(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.6), [
        {"text": "Employee cards loaded live from Microsoft Graph  •  Infinite scroll (100 at a time)  •  Name, title, department, office",
         "size": 13, "color": LIGHT_GRAY, "alignment": PP_ALIGN.LEFT},
    ])

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        "SLIDE 3 — GRID VIEW (show for ~15 seconds)\n"
        "\n"
        "SAY:\n"
        "\"This is the employee directory in Grid View. Each card shows the person's photo, name, "
        "job title, department, and office location — all pulled live from Microsoft Graph.\n"
        "\n"
        "Cards load 100 at a time with infinite scroll, so there's no pagination to deal with. "
        "Just scroll and the next batch loads automatically.\"\n"
        "\n"
        "TRANSITION TO DEMO:\n"
        "After this slide you'll switch to the live browser — have it ready and loaded!\n"
        "\n"
        "TIPS:\n"
        "• This is the slide that should match what the viewer sees in the live demo next.\n"
        "• If you replaced the placeholder with an actual screenshot — briefly point out the cards.\n"
        "• Keep it short — the live demo will show this much better than a static image."
    )

    return slide


def build_slide_4_tech_highlight(prs):
    """Slide 4: Tech Stack Pills
    Modern badge-style layout showing the technology stack
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Top accent bar
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.04), BRAND_BLUE)

    # Title
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.8),
                 "Built on Modern Microsoft 365", font_size=36, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(10), Inches(0.5),
                 "Enterprise-grade architecture with the latest SharePoint Framework",
                 font_size=16, color=LIGHT_GRAY)

    # Accent line
    add_accent_line(slide, Inches(0.8), Inches(2.1), Inches(2.5))

    # Tech stack grid - 2 rows of cards
    card_data = [
        ("SPFx 1.22.2", "Heft Toolchain", "SharePoint Framework with\nthe new Heft build system", BRAND_BLUE),
        ("React 17", "Component Library", "Modern component architecture\nwith hooks and functional components", RGBColor(0x61, 0xDB, 0xFB)),
        ("Fluent UI v9", "Design System", "100% Fluent 2 design tokens\nfor consistent Microsoft look", ACCENT_PURPLE),
        ("Azure OpenAI", "AI Intelligence", "Natural language to OData\nquery translation", RGBColor(0xF2, 0x8A, 0x2E)),
        ("Microsoft Graph", "Data Layer", "Live user profiles, org hierarchy,\nand presence information", SUCCESS_GREEN),
        ("Jotai + IndexedDB", "State & Cache", "Atomic state management with\npersistent local caching", RGBColor(0xD4, 0x73, 0xA8)),
    ]

    x_positions = [Inches(0.8), Inches(4.9), Inches(9.0)]

    for i, (title, subtitle, desc, accent) in enumerate(card_data):
        row = i // 3
        col = i % 3
        x = x_positions[col]
        y = Inches(2.6) + Inches(row * 2.2)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, y, Inches(3.7), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG_ALT
        card.line.color.rgb = RGBColor(0x3A, 0x3A, 0x42)
        card.line.width = Pt(1)
        card.adjustments[0] = 0.06

        # Accent top border on card
        add_shape_rect(slide, x + Pt(6), y + Pt(3), Inches(0.5), Pt(3), accent)

        # Card title
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.3), Inches(0.4),
                     title, font_size=18, color=WHITE, bold=True)

        # Card subtitle
        add_text_box(slide, x + Inches(0.2), y + Inches(0.55), Inches(3.3), Inches(0.3),
                     subtitle, font_size=12, color=accent, bold=True,
                     font_name="Segoe UI Semibold")

        # Card description
        add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.3), Inches(0.7),
                     desc, font_size=11, color=LIGHT_GRAY)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        "SLIDE 4 — TECH STACK (show for ~15-20 seconds)\n"
        "\n"
        "SAY:\n"
        "\"Under the hood, Company Directory is built on the latest Microsoft 365 stack.\n"
        "\n"
        "SharePoint Framework 1.22.2 with the new Heft build toolchain replaces Gulp entirely. "
        "React 17 with functional components and hooks. "
        "Fluent UI version 9 — 100% Fluent 2 design tokens for a native Microsoft look.\n"
        "\n"
        "The AI layer uses Azure OpenAI to translate natural language into OData queries. "
        "Microsoft Graph provides live user profiles and org hierarchy. "
        "And Jotai with IndexedDB handles state management and local caching for fast navigation.\"\n"
        "\n"
        "TIPS:\n"
        "• Don't read every card — hit the highlights: SPFx 1.22.2 (Heft), Azure OpenAI, and Fluent v9.\n"
        "• If short on time, point to 2-3 cards and say \"these are the key building blocks.\"\n"
        "• The judges want to see modern tech choices — emphasize Heft toolchain and AI integration.\n"
        "• This is a good place to mention: \"zero Gulp — fully Heft-based build\" if you want to impress."
    )

    return slide


def build_slide_5_closing_intro(prs):
    """Slide 5: "Let's See It In Action" transition slide
    Minimal, cinematic feel - signals the demo is about to start
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Subtle decorative circles
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2),
                                     Inches(8), Inches(8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x26)
    circle.line.fill.background()

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(2),
                                      Inches(7), Inches(7))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x24)
    circle2.line.fill.background()

    # Center text
    add_text_box(slide, Inches(2), Inches(2.5), Inches(9.333), Inches(1.0),
                 "Let's See It in Action", font_size=44, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Segoe UI Light")

    # Accent line centered
    add_accent_line(slide, Inches(5.5), Inches(3.7), Inches(2.333))

    # Subtitle
    add_text_box(slide, Inches(2), Inches(4.1), Inches(9.333), Inches(0.6),
                 "Live demo on SharePoint Online", font_size=18, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        "SLIDE 5 — TRANSITION TO DEMO (show for ~5 seconds)\n"
        "\n"
        "SAY:\n"
        "\"Now let's see it in action.\"\n"
        "\n"
        "ACTION:\n"
        "• Pause for 2 seconds after saying the line.\n"
        "• Switch to the browser tab with SharePoint already loaded.\n"
        "• Make sure the Grid View is visible with employee cards.\n"
        "• Start the live demo — follow the Video Script from Segment 1 onwards.\n"
        "\n"
        "DEMO ORDER (from HACKATHON_VIDEO_SCRIPT.md):\n"
        "1. Grid View — scroll down to show infinite scroll\n"
        "2. Click 3-dot menu — show chat, email, call actions\n"
        "3. Switch to List View — show data grid with columns\n"
        "4. AI Search — type \"Find all engineers in Porto\"\n"
        "5. AI Search — type \"Managers in the Marketing department\"\n"
        "6. AI Search — type \"Guest users\"\n"
        "7. Org Chart — show top-level, drill into a manager, use UserPicker\n"
        "8. Close — summarize tech stack, multi-host support, thank the viewer\n"
        "\n"
        "TIPS:\n"
        "• Keep the transition smooth — don't fumble with tabs.\n"
        "• Use Cmd+Tab or have the browser as the next window.\n"
        "• Total demo target: ~3 minutes 45 seconds."
    )

    return slide


def build_slide_6_conclusion(prs):
    """Slide 6: Conclusion — Thank You + Social Media Contacts
    Dark background, centered thank-you message, social links, event badge
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Top accent bar
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), BRAND_BLUE)

    # Subtle decorative circles (same style as slide 5)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(-1.5),
                                     Inches(7), Inches(7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x26)
    circle.line.fill.background()
    sp = circle._element
    spTree = sp.getparent()
    spTree.remove(sp)
    spTree.insert(2, sp)

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(2.5),
                                      Inches(6), Inches(6))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x24)
    circle2.line.fill.background()
    sp2 = circle2._element
    spTree2 = sp2.getparent()
    spTree2.remove(sp2)
    spTree2.insert(3, sp2)

    # "Demo Complete" pill badge
    add_pill_badge(slide, Inches(5.1), Inches(1.3), "LIVE DEMO COMPLETE",
                   fill_color=SUCCESS_GREEN, width=Inches(3.2), font_size=13)

    # Main thank-you text
    add_text_box(slide, Inches(2), Inches(2.1), Inches(9.333), Inches(1.0),
                 "Thank You!", font_size=52, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Segoe UI Light")

    # Accent line centered
    add_accent_line(slide, Inches(5.5), Inches(3.3), Inches(2.333))

    # Subtitle
    add_text_box(slide, Inches(2), Inches(3.6), Inches(9.333), Inches(0.6),
                 "Company Directory — AI-Powered People Directory for Microsoft 365",
                 font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Author name
    add_text_box(slide, Inches(2), Inches(4.3), Inches(9.333), Inches(0.5),
                 "João Mendes", font_size=22, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Social media contacts - card style
    social_card_left = Inches(3.5)
    social_card_top = Inches(5.0)
    social_card_width = Inches(6.333)
    social_card_height = Inches(1.6)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   social_card_left, social_card_top,
                                   social_card_width, social_card_height)
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_BG_ALT
    card.line.color.rgb = RGBColor(0x3A, 0x3A, 0x42)
    card.line.width = Pt(1)
    card.adjustments[0] = 0.06

    social_lines = [
        {"text": "\U0001F310  github.com/joaojmendes",
         "size": 14, "color": BRAND_BLUE_LIGHT, "bold": False,
         "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(8)},
        {"text": "\U0001F4BC  linkedin.com/in/intjoaomendes",
         "size": 14, "color": BRAND_BLUE_LIGHT, "bold": False,
         "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(8)},
        {"text": "\U0001D54F  @joaojmendes",
         "size": 14, "color": BRAND_BLUE_LIGHT, "bold": False,
         "alignment": PP_ALIGN.CENTER, "spacing_after": Pt(8)},
    ]
    add_multi_text(slide, social_card_left + Inches(0.3), social_card_top + Inches(0.15),
                   social_card_width - Inches(0.6), social_card_height - Inches(0.3),
                   social_lines)

    # Event badge at bottom
    add_pill_badge(slide, Inches(5.1), Inches(6.9), "SHAREPOINT HACKATHON 2026",
                   fill_color=BRAND_BLUE, width=Inches(3.2), font_size=12)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        "SLIDE 6 — CONCLUSION (show for ~10-15 seconds)\n"
        "\n"
        "SAY:\n"
        "\"And that's Company Directory — an AI-powered people directory built entirely on Microsoft 365.\n"
        "\n"
        "Thank you so much for watching! If you'd like to try it out, the source code is on GitHub. "
        "Feel free to connect with me on LinkedIn or follow me on X for more Microsoft 365 dev content.\n"
        "\n"
        "Thanks again, and happy hacking!\"\n"
        "\n"
        "TIPS:\n"
        "• Smile and speak with energy — leave a positive lasting impression.\n"
        "• Point out the social links on screen so viewers know where to find you.\n"
        "• Keep it brief — 10-15 seconds max. Don't repeat the whole demo.\n"
        "• End confidently — a clear closing is better than trailing off."
    )

    return slide


# ─── Main ────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()

    # Set widescreen 16:9 dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("Building slides...")

    build_slide_1_title(prs)
    print("  ✓ Slide 1: Title Card")

    build_slide_2_overview(prs)
    print("  ✓ Slide 2: Overview + Features")

    build_slide_3_grid_view(prs)
    print("  ✓ Slide 3: Grid View Focus")

    build_slide_4_tech_highlight(prs)
    print("  ✓ Slide 4: Tech Stack")

    build_slide_5_closing_intro(prs)
    print("  ✓ Slide 5: Transition to Demo")

    build_slide_6_conclusion(prs)
    print("  ✓ Slide 6: Conclusion & Thank You")

    prs.save(OUTPUT_PATH)
    print(f"\n✅ Presentation saved to:\n   {OUTPUT_PATH}")
    print(f"\n📝 Notes:")
    print(f"   - Slide 3 has a placeholder for your Grid View screenshot")
    print(f"   - Open in PowerPoint → right-click placeholder → 'Change Picture' to add screenshots")
    print(f"   - All text uses Segoe UI (Fluent 2 default font)")


if __name__ == "__main__":
    main()
